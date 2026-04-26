from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def set_slide_background(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def create_presentation():
    prs = Presentation()

    # Beautiful, vibrant color palette
    colors = {
        "vibrant_purple": RGBColor(100, 30, 150),
        "rich_teal": RGBColor(10, 120, 120),
        "crimson_red": RGBColor(180, 20, 50),
        "deep_indigo": RGBColor(50, 20, 140),
        "sunset_orange": RGBColor(230, 90, 15),
        "emerald_green": RGBColor(15, 140, 75),
        "royal_blue": RGBColor(25, 60, 180),
        "hot_magenta": RGBColor(190, 20, 110),
        "ocean_blue": RGBColor(0, 105, 148),
        "warm_charcoal": RGBColor(50, 50, 60),
        "white": RGBColor(255, 255, 255),
        "bright_yellow": RGBColor(255, 225, 50),
        "neon_cyan": RGBColor(100, 255, 255),
        "soft_peach": RGBColor(255, 210, 170),
        "mint_green": RGBColor(150, 255, 180),
        "dark_text": RGBColor(20, 20, 20)
    }

    slides_data = [
        {
            "layout": 0,
            "title": "StockMind AI\nInstitutional Terminal",
            "subtitle": "AI-Powered Real-Time Stock Market Intelligence\nPresented by Rahul Valluri",
            "bg": "vibrant_purple",
            "title_color": "bright_yellow",
            "text_color": "white"
        },
        {
            "layout": 1,
            "title": "1. Introduction",
            "content": "• StockMind AI is a professional-grade stock market terminal.\n• Delivers real-time market data, AI predictions, and sentiment analysis.\n• Built to bridge the gap between retail investors and institutional trading tools.\n• Features a highly responsive, visually stunning UI.",
            "bg": "rich_teal",
            "title_color": "bright_yellow",
            "text_color": "white"
        },
        {
            "layout": 1,
            "title": "2. Problem Statement",
            "content": "• Retail investors lack access to high-frequency trading (HFT) insights.\n• Existing platforms are cluttered and cause cognitive overload.\n• Disconnected tools: Users need separate apps for news, charts, and AI predictions.\n• High latency in traditional monolithic finance applications.",
            "bg": "crimson_red",
            "title_color": "white",
            "text_color": "soft_peach"
        },
        {
            "layout": 1,
            "title": "3. Proposed Solution",
            "content": "• A unified, all-in-one terminal for market intelligence.\n• Microservices architecture ensuring ultra-low latency updates.\n• Integration of Machine Learning for trend prediction and sentiment tracking.\n• Breathtaking User Interface designed for maximum focus and flow.",
            "bg": "deep_indigo",
            "title_color": "neon_cyan",
            "text_color": "white"
        },
        {
            "layout": 1,
            "title": "4. System Architecture",
            "content": "• Frontend: React + Vite (High-performance UI rendering).\n• Backend Gateway: Node.js + Socket.io (Real-time data streaming).\n• ML Service: Python + FastAPI (Heavy computation & model inference).\n• Authentication: Firebase Auth (Secure user management).",
            "bg": "sunset_orange",
            "title_color": "white",
            "text_color": "dark_text"
        },
        {
            "layout": 1,
            "title": "5. Frontend Technologies",
            "content": "• Built with React, leveraging Lucide React for iconography.\n• Framer Motion used for fluid, cinematic micro-animations.\n• Advanced Recharts integration for complex financial charting.\n• Glassmorphism: Translucent cards with deep volumetric shadows.",
            "bg": "emerald_green",
            "title_color": "bright_yellow",
            "text_color": "white"
        },
        {
            "layout": 1,
            "title": "6. Real-Time Backend",
            "content": "• Node.js handles WebSocket connections for live price ticks.\n• Acts as an API Gateway to the Python ML Service.\n• HFT Simulation: Streams rapid order book updates to the client.\n• Secure architecture using CORS, Helmet, and Compression middleware.",
            "bg": "royal_blue",
            "title_color": "neon_cyan",
            "text_color": "white"
        },
        {
            "layout": 1,
            "title": "7. AI & Machine Learning Engine",
            "content": "• Built entirely in Python using FastAPI for asynchronous speed.\n• Time-Series Forecasting: Utilizes LSTM and Random Forest algorithms.\n• NLP Sentiment Analysis: Extracts financial news via Firecrawl.\n• Algorithmic confidence scoring for bullish/bearish signals.",
            "bg": "hot_magenta",
            "title_color": "bright_yellow",
            "text_color": "white"
        },
        {
            "layout": 1,
            "title": "8. Core Features (Part 1)",
            "content": "• Live Dashboard: Overview of global markets, system health, and latency.\n• Dynamic Market Matrix: Real-time price tracking and volatility heatmaps.\n• Institutional Intelligence: AI parses news and scores market sentiment.\n• Order Book L2: Live visualization of bid/ask spreads.",
            "bg": "ocean_blue",
            "title_color": "mint_green",
            "text_color": "white"
        },
        {
            "layout": 1,
            "title": "9. Core Features (Part 2)",
            "content": "• Predictive Analytics: AI forecasts for short-term and long-term trends.\n• Smart Alerts: Notification system for high-priority market shifts.\n• Portfolio Management: Track assets with real-time valuation updates.\n• Institutional Design: Apple-grade polished aesthetics.",
            "bg": "rich_teal",
            "title_color": "soft_peach",
            "text_color": "white"
        },
        {
            "layout": 0,
            "title": "10. Conclusion & Future Scope",
            "subtitle": "• Expand to real brokerage API integration (Alpaca/Interactive Brokers).\n• Add multi-factor authentication (MFA).\n\nThank you! Any Questions?",
            "bg": "warm_charcoal",
            "title_color": "neon_cyan",
            "text_color": "white"
        }
    ]

    for slide_info in slides_data:
        slide_layout = prs.slide_layouts[slide_info["layout"]]
        slide = prs.slides.add_slide(slide_layout)
        
        set_slide_background(slide, colors[slide_info["bg"]])

        title_shape = slide.shapes.title
        title_shape.text = slide_info["title"]
        
        for paragraph in title_shape.text_frame.paragraphs:
            paragraph.font.color.rgb = colors[slide_info["title_color"]]
            paragraph.font.bold = True
            paragraph.alignment = PP_ALIGN.CENTER if slide_info["layout"] == 0 else PP_ALIGN.LEFT

        if slide_info["layout"] == 0:
            subtitle_shape = slide.placeholders[1]
            subtitle_shape.text = slide_info["subtitle"]
            for paragraph in subtitle_shape.text_frame.paragraphs:
                paragraph.font.color.rgb = colors[slide_info["text_color"]]
                paragraph.alignment = PP_ALIGN.CENTER
        else:
            content_shape = slide.placeholders[1]
            content_shape.text = slide_info["content"]
            for paragraph in content_shape.text_frame.paragraphs:
                paragraph.font.color.rgb = colors[slide_info["text_color"]]
                paragraph.font.size = Pt(24)

    prs.save('StockMind_AI_Colorful_Presentation.pptx')
    print("Successfully generated StockMind_AI_Colorful_Presentation.pptx")

if __name__ == '__main__':
    create_presentation()
